"""Service for generating professional PowerPoint presentations for RFP soutenance."""
import io
import logging
import math
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

# ── Color palette ──
COLOR_PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)
COLOR_SECONDARY = RGBColor(0x2C, 0x5F, 0x8A)
COLOR_ACCENT = RGBColor(0x3D, 0x7A, 0xB5)
COLOR_ACCENT_LIGHT = RGBColor(0x5A, 0x9F, 0xD4)
COLOR_ACCENT_PALE = RGBColor(0xD6, 0xEA, 0xF8)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
COLOR_DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
COLOR_MUTED = RGBColor(0x7F, 0x8C, 0x8D)
COLOR_SUCCESS = RGBColor(0x27, 0xAE, 0x60)
COLOR_SUCCESS_LIGHT = RGBColor(0xD5, 0xF5, 0xE3)
COLOR_WARN = RGBColor(0xE6, 0x7E, 0x22)
COLOR_WARN_LIGHT = RGBColor(0xFD, 0xEA, 0xD3)
COLOR_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
COLOR_PURPLE_LIGHT = RGBColor(0xEB, 0xDE, 0xF0)
COLOR_TEAL = RGBColor(0x00, 0x89, 0x7B)
COLOR_TEAL_LIGHT = RGBColor(0xD1, 0xF2, 0xEB)
COLOR_CORAL = RGBColor(0xE7, 0x4C, 0x3C)
COLOR_CORAL_LIGHT = RGBColor(0xFA, 0xDB, 0xD8)
COLOR_SIDEBAR = RGBColor(0xE8, 0xEE, 0xF5)

# Unicode icons for visual decoration
ICON_CHECK = "\u2713"       # ✓
ICON_ARROW_R = "\u279C"     # ➜
ICON_STAR = "\u2605"        # ★
ICON_DIAMOND = "\u25C6"     # ◆
ICON_CIRCLE = "\u25CF"      # ●
ICON_BULLET = "\u2022"      # •
ICON_TARGET = "\u25CE"      # ◎
ICON_SQUARE = "\u25A0"      # ■
ICON_TRIANGLE = "\u25B6"    # ▶
ICON_RING = "\u25CB"        # ○
ICON_DASH = "\u2014"        # —
ICON_CHEVRON_R = "\u276F"  # ❯
ICON_QUOTE_L = "\u201C"    # "
ICON_QUOTE_R = "\u201D"    # "
ICON_PERSON = "\u2B24"     # ⬤ (person placeholder)
ICON_LIGHT = "\U0001F4A1"  # 💡
ICON_ROCKET = "\U0001F680" # 🚀
ICON_GEAR = "\u2699"       # ⚙
ICON_CHART = "\U0001F4CA"  # 📊
ICON_SHIELD = "\U0001F6E1" # 🛡
ICON_HANDSHAKE = "\U0001F91D" # 🤝
ICON_TROPHY = "\U0001F3C6"   # 🏆
ICON_CALENDAR = "\U0001F4C5" # 📅
ICON_PEOPLE = "\U0001F465"   # 👥
ICON_TOOLS = "\U0001F527"    # 🔧

# Layout type constants
LAYOUT_BULLETS = "bullets"
LAYOUT_PROCESS = "process_flow"
LAYOUT_TIMELINE = "timeline"
LAYOUT_ICON_GRID = "icon_grid"
LAYOUT_COMPARISON = "comparison"
LAYOUT_QUOTE = "quote"
LAYOUT_PYRAMID = "pyramid"
LAYOUT_DASHBOARD = "dashboard"

# Section-specific icons (mapped by keyword detection)
SECTION_ICONS = {
    "contexte": "\U0001F50D",     # 🔍
    "comprehen": "\U0001F4CB",    # 📋
    "solution": "\U0001F4A1",     # 💡
    "methodol": "\u2699",         # ⚙
    "approche": "\U0001F4CA",     # 📊
    "equipe": "\U0001F465",       # 👥
    "moyen": "\U0001F527",        # 🔧
    "planning": "\U0001F4C5",     # 📅
    "livrable": "\U0001F4E6",     # 📦
    "reference": "\U0001F3C6",    # 🏆
    "experience": "\U0001F3C6",   # 🏆
    "qualite": "\u2705",          # ✅
    "engage": "\U0001F91D",       # 🤝
    "valeur": "\U0001F4B0",       # 💰
    "innov": "\U0001F680",        # 🚀
    "securite": "\U0001F6E1",     # 🛡
    "gouvern": "\U0001F3DB",      # 🏛
    "transition": "\U0001F504",   # 🔄
    "support": "\U0001F6E0",      # 🛠
    "form": "\U0001F393",         # 🎓
}

# Color rotation for sections
SECTION_COLORS = [
    (COLOR_SECONDARY, COLOR_ACCENT_PALE),
    (COLOR_TEAL, COLOR_TEAL_LIGHT),
    (COLOR_PURPLE, COLOR_PURPLE_LIGHT),
    (COLOR_SUCCESS, COLOR_SUCCESS_LIGHT),
    (COLOR_WARN, COLOR_WARN_LIGHT),
    (COLOR_CORAL, COLOR_CORAL_LIGHT),
    (COLOR_ACCENT, COLOR_ACCENT_PALE),
    (COLOR_PRIMARY, COLOR_LIGHT_BG),
]


def _get_section_icon(title: str) -> str:
    """Get an appropriate icon for a section based on its title."""
    lower = title.lower()
    for keyword, icon in SECTION_ICONS.items():
        if keyword in lower:
            return icon
    return ICON_DIAMOND


def _normalize_bullet(item) -> str:
    """Ensure a bullet point is a plain string, regardless of AI output format."""
    if isinstance(item, str):
        return item
    if isinstance(item, dict):
        # AI sometimes returns {"text": "...", "icon": "..."} or {"value": "...", "label": "..."}
        return (
            item.get("text")
            or item.get("value")
            or item.get("label")
            or item.get("title")
            or item.get("content")
            or item.get("description")
            or str(item)
        )
    return str(item)


class RFPPptxService:
    """Service for generating professional soutenance PowerPoint presentations."""

    SLIDE_W = Inches(13.333)  # 16:9 widescreen
    SLIDE_H = Inches(7.5)

    # ── Low-level helpers ──

    @staticmethod
    def _add_background(slide, color: RGBColor):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    @staticmethod
    def _rect(slide, left, top, width, height, color: RGBColor):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    @staticmethod
    def _rounded_rect(slide, left, top, width, height, color: RGBColor,
                      border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    @staticmethod
    def _oval(slide, left, top, w, h, color: RGBColor):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    @classmethod
    def _icon_badge(cls, slide, left, top, size, bg_color, icon_text,
                    icon_size=16, icon_color=None):
        """Create a colored circle with an icon/text inside."""
        circ = cls._oval(slide, left, top, size, size, bg_color)
        tf = circ.text_frame
        tf.paragraphs[0].text = icon_text
        tf.paragraphs[0].font.size = Pt(icon_size)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = icon_color or COLOR_WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        return circ

    @staticmethod
    def _text_box(slide, left, top, width, height, text: str,
                  font_size=14, bold=False, color=COLOR_DARK_TEXT,
                  alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    @staticmethod
    def _add_notes(slide, text: str):
        if text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = text

    # ── Decorative/layout elements ──

    @classmethod
    def _add_footer_bar(cls, slide, text="CONFIDENTIEL", section_color=None):
        """Professional footer with accent line."""
        cls._rect(slide, Inches(0), cls.SLIDE_H - Inches(0.4),
                  cls.SLIDE_W, Inches(0.4), COLOR_PRIMARY)
        cls._rect(slide, Inches(0), cls.SLIDE_H - Inches(0.43),
                  cls.SLIDE_W, Inches(0.03), section_color or COLOR_ACCENT)
        if text:
            cls._text_box(slide, Inches(0.5), cls.SLIDE_H - Inches(0.37),
                          Inches(6), Inches(0.32), text,
                          font_size=8, color=COLOR_WHITE)

    @classmethod
    def _add_left_sidebar(cls, slide, accent_color=COLOR_ACCENT, icon=""):
        """Add a thin decorative sidebar on the left."""
        # Sidebar strip
        cls._rect(slide, Inches(0), Inches(0), Inches(0.45), cls.SLIDE_H, COLOR_SIDEBAR)
        # Accent color bar
        cls._rect(slide, Inches(0), Inches(0), Inches(0.06), cls.SLIDE_H, accent_color)
        # Icon in sidebar
        if icon:
            cls._text_box(slide, Inches(0.06), Inches(3.2), Inches(0.38), Inches(0.4),
                          icon, font_size=16, color=accent_color, alignment=PP_ALIGN.CENTER)

    @classmethod
    def _add_slide_header(cls, slide, title: str, subtitle: str = "",
                          section_label: str = "", accent_color=None):
        """Professional slide header bar."""
        ac = accent_color or COLOR_ACCENT
        # Top accent line
        cls._rect(slide, Inches(0), Inches(0), cls.SLIDE_W, Inches(0.05), ac)
        # Header background
        cls._rect(slide, Inches(0), Inches(0.05), cls.SLIDE_W, Inches(1.0), COLOR_PRIMARY)

        x_offset = Inches(0.8)

        if section_label:
            cls._text_box(slide, x_offset, Inches(0.12), Inches(10), Inches(0.3),
                          section_label.upper(), font_size=9, bold=True,
                          color=ac)

        title_y = Inches(0.35) if section_label else Inches(0.2)
        cls._text_box(slide, x_offset, title_y, Inches(11), Inches(0.6),
                      title, font_size=24, bold=True, color=COLOR_WHITE)

        if subtitle:
            cls._text_box(slide, x_offset, Inches(1.15), Inches(10), Inches(0.35),
                          subtitle, font_size=13, color=COLOR_MUTED)

    @classmethod
    def _add_decorative_dots(cls, slide, x_start, y_start, cols=5, rows=4,
                             color=COLOR_ACCENT_PALE, dot_size=Inches(0.08),
                             spacing=Inches(0.25)):
        """Add a grid of small decorative dots (modern pattern)."""
        for r in range(rows):
            for c in range(cols):
                cls._oval(slide, x_start + c * spacing, y_start + r * spacing,
                          dot_size, dot_size, color)

    # ── Slide types ──

    @classmethod
    def _create_title_slide(cls, prs, project_name, client_name, company_name, rfp_reference):
        """Cover slide with rich decorative geometry and modern layout."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_PRIMARY)

        # Large decorative circles (overlapping, modern depth effect)
        cls._oval(slide, cls.SLIDE_W - Inches(6), Inches(-3),
                  Inches(10), Inches(10), COLOR_SECONDARY)
        cls._oval(slide, cls.SLIDE_W - Inches(4), Inches(-1.5),
                  Inches(6), Inches(6), COLOR_ACCENT)
        cls._oval(slide, Inches(-4), cls.SLIDE_H - Inches(5),
                  Inches(8), Inches(8), COLOR_SECONDARY)
        # Additional subtle circles for depth
        cls._oval(slide, Inches(6), Inches(-4),
                  Inches(5), Inches(5), COLOR_SECONDARY)
        cls._oval(slide, cls.SLIDE_W - Inches(2), Inches(4),
                  Inches(3), Inches(3), COLOR_ACCENT)

        # Dot pattern decorations (two clusters)
        cls._add_decorative_dots(slide, Inches(0.5), Inches(0.3), cols=6, rows=3,
                                 color=COLOR_ACCENT_LIGHT, spacing=Inches(0.2))
        cls._add_decorative_dots(slide, cls.SLIDE_W - Inches(3), cls.SLIDE_H - Inches(2),
                                 cols=5, rows=4, color=COLOR_ACCENT_LIGHT, spacing=Inches(0.2))

        # Top accent bar (dual-tone)
        cls._rect(slide, Inches(0), Inches(0), cls.SLIDE_W, Inches(0.08), COLOR_ACCENT)
        cls._rect(slide, Inches(0), Inches(0.08), cls.SLIDE_W, Inches(0.03), COLOR_TEAL)

        # Left vertical accent bar
        cls._rect(slide, Inches(0.65), Inches(1.2), Inches(0.06), Inches(4.5), COLOR_ACCENT)

        # "SOUTENANCE" label with icon
        cls._text_box(slide, Inches(1), Inches(1.3), Inches(8), Inches(0.5),
                      f"{ICON_TROPHY}  SOUTENANCE COMMERCIALE", font_size=14, bold=True,
                      color=COLOR_ACCENT_LIGHT)

        # Project name (larger, more prominent)
        cls._text_box(slide, Inches(1), Inches(2.0), Inches(9), Inches(1.8),
                      project_name, font_size=42, bold=True, color=COLOR_WHITE)

        # Double separator lines
        cls._rect(slide, Inches(1), Inches(3.9), Inches(4), Inches(0.05), COLOR_ACCENT)
        cls._rect(slide, Inches(1), Inches(4.0), Inches(2.5), Inches(0.03), COLOR_TEAL)

        # Info block with richer design
        info_y = Inches(4.4)
        if client_name:
            cls._icon_badge(slide, Inches(1), info_y, Inches(0.4), COLOR_ACCENT,
                            ICON_TARGET, icon_size=14)
            cls._text_box(slide, Inches(1.55), info_y + Inches(0.02), Inches(7), Inches(0.4),
                          f"Client : {client_name}", font_size=17, bold=True,
                          color=COLOR_WHITE)
            info_y += Inches(0.6)
        if company_name:
            cls._icon_badge(slide, Inches(1), info_y, Inches(0.4), COLOR_TEAL,
                            ICON_DIAMOND, icon_size=14)
            cls._text_box(slide, Inches(1.55), info_y + Inches(0.02), Inches(7), Inches(0.4),
                          f"Soumissionnaire : {company_name}", font_size=17, bold=True,
                          color=COLOR_WHITE)
            info_y += Inches(0.6)
        if rfp_reference:
            cls._icon_badge(slide, Inches(1), info_y, Inches(0.4), COLOR_PURPLE,
                            ICON_SQUARE, icon_size=12)
            cls._text_box(slide, Inches(1.55), info_y + Inches(0.02), Inches(7), Inches(0.4),
                          f"Reference : {rfp_reference}", font_size=14, color=COLOR_ACCENT_LIGHT)

        # Bottom bar (dual-tone)
        cls._rect(slide, Inches(0), cls.SLIDE_H - Inches(0.55),
                  cls.SLIDE_W, Inches(0.55), COLOR_SECONDARY)
        cls._rect(slide, Inches(0), cls.SLIDE_H - Inches(0.55),
                  cls.SLIDE_W, Inches(0.04), COLOR_ACCENT)
        cls._text_box(slide, Inches(1), cls.SLIDE_H - Inches(0.45),
                      Inches(10), Inches(0.4),
                      "DOCUMENT CONFIDENTIEL", font_size=10, bold=True,
                      color=COLOR_WHITE)

    @classmethod
    def _create_agenda_slide(cls, prs, sections):
        """Agenda slide with visual timeline and colored dots."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_WHITE)
        cls._add_slide_header(slide, "AGENDA")

        # Decorative dots in bottom-right
        cls._add_decorative_dots(slide, cls.SLIDE_W - Inches(2), cls.SLIDE_H - Inches(2),
                                 cols=4, rows=4, color=COLOR_ACCENT_PALE)

        # Timeline line
        line_x = Inches(1.6)
        cls._rect(slide, line_x + Inches(0.19), Inches(1.5),
                  Inches(0.04), Inches(5.3), COLOR_ACCENT_PALE)

        y = Inches(1.5)
        for i, section in enumerate(sections):
            title = section.get("title", "")
            duration = section.get("duration", "")
            sc, sc_light = SECTION_COLORS[i % len(SECTION_COLORS)]
            icon = _get_section_icon(title)

            # Timeline dot with icon
            cls._icon_badge(slide, line_x, y + Inches(0.03),
                            Inches(0.42), sc, icon, icon_size=14)

            # Title
            cls._text_box(slide, Inches(2.3), y, Inches(7), Inches(0.45),
                          title, font_size=15, bold=True, color=COLOR_PRIMARY)

            # Duration badge
            if duration:
                cls._rounded_rect(slide, Inches(10), y + Inches(0.05),
                                  Inches(2), Inches(0.35), sc_light, border_color=sc)
                cls._text_box(slide, Inches(10), y + Inches(0.05),
                              Inches(2), Inches(0.35), f"{ICON_CIRCLE} {duration}",
                              font_size=10, color=sc, alignment=PP_ALIGN.CENTER)

            y += Inches(0.58)

        cls._add_footer_bar(slide)

    @classmethod
    def _create_section_divider(cls, prs, section_number, section_title,
                                duration="", total_sections=1):
        """Visually striking section divider with rich geometric decoration."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sc, sc_light = SECTION_COLORS[(section_number - 1) % len(SECTION_COLORS)]

        cls._add_background(slide, COLOR_SECONDARY)

        # Large decorative shapes - more layers for depth
        cls._oval(slide, Inches(-3), Inches(-3),
                  Inches(8), Inches(8), COLOR_PRIMARY)
        cls._oval(slide, cls.SLIDE_W - Inches(5), cls.SLIDE_H - Inches(5),
                  Inches(8), Inches(8), sc)
        cls._oval(slide, cls.SLIDE_W - Inches(2.5), Inches(0),
                  Inches(2), Inches(2), COLOR_ACCENT)
        cls._oval(slide, Inches(4), Inches(-1.5),
                  Inches(3), Inches(3), COLOR_PRIMARY)

        # Dot patterns (two clusters)
        cls._add_decorative_dots(slide, Inches(8), Inches(0.5), cols=5, rows=3,
                                 color=COLOR_ACCENT_LIGHT, spacing=Inches(0.3))
        cls._add_decorative_dots(slide, Inches(1), cls.SLIDE_H - Inches(2),
                                 cols=4, rows=3, color=sc_light, spacing=Inches(0.25))

        # Section icon (large, as background element)
        icon = _get_section_icon(section_title)
        cls._text_box(slide, cls.SLIDE_W - Inches(3.5), Inches(0.3),
                      Inches(3), Inches(2), icon, font_size=60,
                      color=sc_light, alignment=PP_ALIGN.CENTER)

        # Section number with background circle
        num_text = f"0{section_number}" if section_number < 10 else str(section_number)
        cls._icon_badge(slide, Inches(5.6), Inches(1.0),
                        Inches(2.1), sc, num_text, icon_size=50)

        # Section title
        cls._text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                      section_title.upper(), font_size=32, bold=True,
                      color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

        # Double accent bars
        cls._rect(slide, Inches(4), Inches(4.6), Inches(5.3), Inches(0.05), sc)
        cls._rect(slide, Inches(4.5), Inches(4.7), Inches(4.3), Inches(0.03),
                  COLOR_ACCENT_LIGHT)

        # Duration + icon badge
        if duration:
            cls._rounded_rect(slide, Inches(4.8), Inches(5.0),
                              Inches(3.6), Inches(0.5), sc_light)
            cls._text_box(slide, Inches(4.8), Inches(5.02),
                          Inches(3.6), Inches(0.45),
                          f"{icon}  {duration}", font_size=16, bold=True,
                          color=sc, alignment=PP_ALIGN.CENTER)

        # Progress bar at bottom (segmented style)
        bar_y = cls.SLIDE_H - Inches(0.35)
        cls._rect(slide, Inches(0), bar_y, cls.SLIDE_W, Inches(0.35), COLOR_PRIMARY)
        if total_sections > 0:
            progress_w = int(cls.SLIDE_W) * section_number // total_sections
            cls._rect(slide, Inches(0), bar_y, progress_w, Inches(0.35), sc)
            # Progress indicator dot
            cls._oval(slide, progress_w - Inches(0.15), bar_y + Inches(0.05),
                      Inches(0.25), Inches(0.25), COLOR_WHITE)

    @classmethod
    def _create_content_slide(cls, prs, title, bullet_points, subtitle="",
                              speaker_notes="", section_label="",
                              section_idx=0):
        """Content slide with sidebar accent and icon bullets."""
        bullet_points = [_normalize_bullet(b) for b in bullet_points]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_WHITE)

        sc, sc_light = SECTION_COLORS[section_idx % len(SECTION_COLORS)]
        icon = _get_section_icon(section_label or title)

        # Left sidebar
        cls._add_left_sidebar(slide, accent_color=sc, icon=icon)

        cls._add_slide_header(slide, title, subtitle, section_label, accent_color=sc)
        cls._add_footer_bar(slide, section_color=sc)

        start_y = Inches(1.5) if subtitle else Inches(1.3)
        content_x = Inches(0.65)
        content_w = Inches(12.2)

        n_bullets = len(bullet_points)

        if n_bullets >= 6:
            # Two-column layout
            mid = (n_bullets + 1) // 2
            cls._add_icon_bullet_list(slide, content_x, start_y,
                                      Inches(5.8), Inches(5),
                                      bullet_points[:mid], sc, sc_light)
            cls._rect(slide, Inches(6.6), start_y + Inches(0.2),
                      Inches(0.02), Inches(4.5), COLOR_LIGHT_BG)
            cls._add_icon_bullet_list(slide, Inches(6.9), start_y,
                                      Inches(5.8), Inches(5),
                                      bullet_points[mid:], sc, sc_light)
        elif n_bullets <= 4 and all(len(b) < 100 for b in bullet_points):
            # Card layout
            cls._add_card_layout(slide, start_y, bullet_points, sc, sc_light)
        else:
            # Standard bullet list with icons
            cls._add_icon_bullet_list(slide, content_x, start_y,
                                      content_w, Inches(5.2),
                                      bullet_points, sc, sc_light)

        # Decorative dots in bottom-right corner
        cls._add_decorative_dots(slide, cls.SLIDE_W - Inches(1.5),
                                 cls.SLIDE_H - Inches(1.5),
                                 cols=3, rows=3, color=sc_light)

        cls._add_notes(slide, speaker_notes)

    @classmethod
    def _add_icon_bullet_list(cls, slide, left, top, width, height,
                              bullets, accent_color, light_color):
        """Bullet list where each item has a small colored icon badge."""
        y = top
        for idx, point in enumerate(bullets):
            clean = point.lstrip("- ").strip()
            is_sub = point.strip().startswith("- ") and point.startswith("  ")

            if is_sub:
                cls._text_box(slide, left + Inches(0.7), y, width - Inches(0.8), Inches(0.4),
                              f"{ICON_ARROW_R}  {clean}", font_size=12,
                              color=COLOR_MUTED)
                y += Inches(0.38)
            else:
                # Icon dot for main bullets
                cls._icon_badge(slide, left, y + Inches(0.02),
                                Inches(0.28), accent_color,
                                ICON_TRIANGLE, icon_size=9)
                cls._text_box(slide, left + Inches(0.4), y, width - Inches(0.5), Inches(0.5),
                              clean, font_size=14, color=COLOR_DARK_TEXT)
                y += Inches(0.5)

    @classmethod
    def _add_card_layout(cls, slide, start_y, items, accent_color, light_color):
        """Card-style layout with icon badges."""
        cols = 2
        col_width = Inches(5.8)
        card_h = Inches(2.0)
        x_start = Inches(0.8)
        gap_x = Inches(0.5)
        gap_y = Inches(0.4)

        icons = [ICON_STAR, ICON_TARGET, ICON_DIAMOND, ICON_CHECK, ICON_CIRCLE, ICON_SQUARE]
        colors_cycle = [COLOR_SECONDARY, COLOR_TEAL, COLOR_PURPLE,
                        COLOR_SUCCESS, COLOR_WARN, COLOR_ACCENT]

        for i, item in enumerate(items):
            col = i % cols
            row = i // cols
            x = x_start + col * (col_width + gap_x)
            y = start_y + row * (card_h + gap_y)
            ic = colors_cycle[i % len(colors_cycle)]

            # Card background with colored left border
            cls._rounded_rect(slide, x, y, col_width, card_h, COLOR_CARD_BG)
            cls._rect(slide, x, y + Inches(0.1), Inches(0.05), card_h - Inches(0.2), ic)

            # Icon badge
            cls._icon_badge(slide, x + Inches(0.3), y + Inches(0.35),
                            Inches(0.55), ic, icons[i % len(icons)], icon_size=18)

            # Text
            clean = item.lstrip("- ").strip()
            cls._text_box(slide, x + Inches(1.05), y + Inches(0.25),
                          col_width - Inches(1.4), card_h - Inches(0.5),
                          clean, font_size=13, color=COLOR_DARK_TEXT)

    # ── Visual layout slides ──

    @classmethod
    def _create_process_flow_slide(cls, prs, title, steps, subtitle="",
                                    speaker_notes="", section_label="",
                                    section_idx=0):
        """Process flow slide with connected chevron/arrow steps."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_WHITE)

        sc, sc_light = SECTION_COLORS[section_idx % len(SECTION_COLORS)]
        icon = _get_section_icon(section_label or title)

        cls._add_left_sidebar(slide, accent_color=sc, icon=icon)
        cls._add_slide_header(slide, title, subtitle, section_label, accent_color=sc)
        cls._add_footer_bar(slide, section_color=sc)

        n_steps = min(len(steps), 6)
        if n_steps == 0:
            cls._add_notes(slide, speaker_notes)
            return

        # Calculate layout
        total_w = Inches(11.8)
        start_x = Inches(0.7)
        step_w = total_w / n_steps
        step_h = Inches(3.8)
        step_y = Inches(1.6)

        step_icons = [ICON_TARGET, ICON_GEAR, ICON_TOOLS, ICON_CHART,
                      ICON_ROCKET, ICON_TROPHY, ICON_SHIELD, ICON_HANDSHAKE]

        for i, step_text in enumerate(steps[:n_steps]):
            clean = step_text.lstrip("- ").strip()
            x = start_x + i * step_w
            item_sc, item_sc_light = SECTION_COLORS[(section_idx + i) % len(SECTION_COLORS)]

            # Step card
            card_w = step_w - Inches(0.15)
            card_h = step_h

            # Rounded card background
            cls._rounded_rect(slide, x, step_y, card_w, card_h, item_sc_light)
            # Top colored bar on card
            cls._rect(slide, x + Inches(0.05), step_y, card_w - Inches(0.1),
                      Inches(0.06), item_sc)

            # Step number circle
            num_size = Inches(0.55)
            num_x = x + (card_w - num_size) / 2
            cls._icon_badge(slide, num_x, step_y + Inches(0.25),
                            num_size, item_sc, str(i + 1), icon_size=18)

            # Step icon below number
            icon_text = step_icons[i % len(step_icons)]
            cls._text_box(slide, x, step_y + Inches(0.95), card_w, Inches(0.4),
                          icon_text, font_size=20, color=item_sc,
                          alignment=PP_ALIGN.CENTER)

            # Step text
            cls._text_box(slide, x + Inches(0.1), step_y + Inches(1.4),
                          card_w - Inches(0.2), card_h - Inches(1.7),
                          clean, font_size=11, color=COLOR_DARK_TEXT,
                          alignment=PP_ALIGN.CENTER)

            # Arrow connector between steps
            if i < n_steps - 1:
                arrow_x = x + card_w - Inches(0.05)
                cls._text_box(slide, arrow_x, step_y + Inches(0.3),
                              Inches(0.3), Inches(0.4),
                              ICON_CHEVRON_R, font_size=18, bold=True,
                              color=COLOR_MUTED, alignment=PP_ALIGN.CENTER)

        cls._add_notes(slide, speaker_notes)

    @classmethod
    def _create_timeline_slide(cls, prs, title, milestones, subtitle="",
                                speaker_notes="", section_label="",
                                section_idx=0):
        """Timeline slide with connected milestones on a horizontal line."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_WHITE)

        sc, sc_light = SECTION_COLORS[section_idx % len(SECTION_COLORS)]
        icon = _get_section_icon(section_label or title)

        cls._add_left_sidebar(slide, accent_color=sc, icon=icon)
        cls._add_slide_header(slide, title, subtitle, section_label, accent_color=sc)
        cls._add_footer_bar(slide, section_color=sc)

        n_items = min(len(milestones), 8)
        if n_items == 0:
            cls._add_notes(slide, speaker_notes)
            return

        # Horizontal timeline line
        line_y = Inches(3.6)
        line_start_x = Inches(0.8)
        line_end_x = Inches(12.5)
        line_w = line_end_x - line_start_x

        # Gradient-like timeline bar (layered rectangles)
        cls._rect(slide, line_start_x, line_y, line_w, Inches(0.06), sc)
        # Glow effect
        cls._rect(slide, line_start_x, line_y - Inches(0.02), line_w,
                  Inches(0.1), sc_light)
        cls._rect(slide, line_start_x, line_y, line_w, Inches(0.06), sc)

        spacing = line_w / (n_items)

        for i, item_text in enumerate(milestones[:n_items]):
            clean = item_text.lstrip("- ").strip()
            item_sc, item_sc_light = SECTION_COLORS[(section_idx + i) % len(SECTION_COLORS)]
            x_center = line_start_x + spacing * (i + 0.5)

            # Milestone dot on timeline
            dot_size = Inches(0.35)
            cls._oval(slide, x_center - dot_size / 2, line_y - dot_size / 2 + Inches(0.03),
                      dot_size, dot_size, item_sc)
            # Inner white dot
            inner = Inches(0.15)
            cls._oval(slide, x_center - inner / 2, line_y - inner / 2 + Inches(0.03),
                      inner, inner, COLOR_WHITE)

            # Alternate above/below for readability
            text_w = spacing - Inches(0.1)
            if i % 2 == 0:
                # Above the line
                # Small connector line
                cls._rect(slide, x_center - Inches(0.01), line_y - Inches(0.8),
                          Inches(0.02), Inches(0.65), item_sc_light)
                # Label card
                card_y = Inches(1.5)
                card_h = Inches(1.8)
                cls._rounded_rect(slide, x_center - text_w / 2, card_y,
                                  text_w, card_h, item_sc_light)
                cls._rect(slide, x_center - text_w / 2, card_y + card_h - Inches(0.05),
                          text_w, Inches(0.05), item_sc)
                # Step number
                cls._text_box(slide, x_center - text_w / 2, card_y + Inches(0.08),
                              text_w, Inches(0.3),
                              f"ETAPE {i + 1}", font_size=8, bold=True,
                              color=item_sc, alignment=PP_ALIGN.CENTER)
                cls._text_box(slide, x_center - text_w / 2 + Inches(0.05),
                              card_y + Inches(0.35),
                              text_w - Inches(0.1), card_h - Inches(0.5),
                              clean, font_size=10, color=COLOR_DARK_TEXT,
                              alignment=PP_ALIGN.CENTER)
            else:
                # Below the line
                cls._rect(slide, x_center - Inches(0.01), line_y + Inches(0.15),
                          Inches(0.02), Inches(0.65), item_sc_light)
                card_y = line_y + Inches(0.9)
                card_h = Inches(1.8)
                cls._rounded_rect(slide, x_center - text_w / 2, card_y,
                                  text_w, card_h, item_sc_light)
                cls._rect(slide, x_center - text_w / 2, card_y,
                          text_w, Inches(0.05), item_sc)
                cls._text_box(slide, x_center - text_w / 2, card_y + Inches(0.08),
                              text_w, Inches(0.3),
                              f"ETAPE {i + 1}", font_size=8, bold=True,
                              color=item_sc, alignment=PP_ALIGN.CENTER)
                cls._text_box(slide, x_center - text_w / 2 + Inches(0.05),
                              card_y + Inches(0.35),
                              text_w - Inches(0.1), card_h - Inches(0.5),
                              clean, font_size=10, color=COLOR_DARK_TEXT,
                              alignment=PP_ALIGN.CENTER)

        cls._add_notes(slide, speaker_notes)

    @classmethod
    def _create_icon_grid_slide(cls, prs, title, items, subtitle="",
                                 speaker_notes="", section_label="",
                                 section_idx=0):
        """Grid of icon cards - great for features, capabilities, values."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_WHITE)

        sc, sc_light = SECTION_COLORS[section_idx % len(SECTION_COLORS)]
        icon = _get_section_icon(section_label or title)

        cls._add_left_sidebar(slide, accent_color=sc, icon=icon)
        cls._add_slide_header(slide, title, subtitle, section_label, accent_color=sc)
        cls._add_footer_bar(slide, section_color=sc)

        n_items = min(len(items), 9)
        if n_items == 0:
            cls._add_notes(slide, speaker_notes)
            return

        # Grid layout: determine rows and cols
        if n_items <= 3:
            cols, rows = n_items, 1
        elif n_items <= 6:
            cols = 3
            rows = (n_items + 2) // 3
        else:
            cols = 3
            rows = 3

        grid_w = Inches(11.5)
        grid_h = Inches(5.0)
        card_w = grid_w / cols - Inches(0.3)
        card_h = grid_h / rows - Inches(0.3)
        start_x = Inches(0.9)
        start_y = Inches(1.5)

        grid_icons = [ICON_ROCKET, ICON_LIGHT, ICON_GEAR, ICON_SHIELD,
                      ICON_CHART, ICON_TROPHY, ICON_HANDSHAKE, ICON_TOOLS, ICON_CALENDAR]

        for i, item_text in enumerate(items[:n_items]):
            if i >= cols * rows:
                break
            clean = item_text.lstrip("- ").strip()
            col = i % cols
            row = i // cols
            item_sc, item_sc_light = SECTION_COLORS[(section_idx + i) % len(SECTION_COLORS)]

            x = start_x + col * (card_w + Inches(0.3))
            y = start_y + row * (card_h + Inches(0.3))

            # Card background with subtle border
            cls._rounded_rect(slide, x, y, card_w, card_h, COLOR_CARD_BG,
                              border_color=item_sc_light)
            # Top colored accent bar
            cls._rect(slide, x + Inches(0.15), y, card_w - Inches(0.3),
                      Inches(0.05), item_sc)

            # Large icon
            icon_text = grid_icons[i % len(grid_icons)]
            icon_size = Inches(0.65)
            cls._icon_badge(slide, x + (card_w - icon_size) / 2,
                            y + Inches(0.2), icon_size, item_sc,
                            icon_text, icon_size=22)

            # Text below icon
            cls._text_box(slide, x + Inches(0.15), y + Inches(1.0),
                          card_w - Inches(0.3), card_h - Inches(1.15),
                          clean, font_size=12, color=COLOR_DARK_TEXT,
                          alignment=PP_ALIGN.CENTER)

        cls._add_notes(slide, speaker_notes)

    @classmethod
    def _create_comparison_slide(cls, prs, title, items, subtitle="",
                                  speaker_notes="", section_label="",
                                  section_idx=0):
        """Side-by-side comparison or before/after slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_WHITE)

        sc, sc_light = SECTION_COLORS[section_idx % len(SECTION_COLORS)]
        icon = _get_section_icon(section_label or title)

        cls._add_left_sidebar(slide, accent_color=sc, icon=icon)
        cls._add_slide_header(slide, title, subtitle, section_label, accent_color=sc)
        cls._add_footer_bar(slide, section_color=sc)

        n_items = len(items)
        if n_items == 0:
            cls._add_notes(slide, speaker_notes)
            return

        # Split items into two columns for comparison
        mid = (n_items + 1) // 2
        left_items = items[:mid]
        right_items = items[mid:]

        col_w = Inches(5.6)
        col_h = Inches(4.8)
        left_x = Inches(0.7)
        right_x = Inches(6.8)
        col_y = Inches(1.5)

        # Left column - larger card
        cls._rounded_rect(slide, left_x, col_y, col_w, col_h,
                          sc_light, border_color=sc)
        cls._rect(slide, left_x, col_y, col_w, Inches(0.5), sc)
        cls._text_box(slide, left_x, col_y + Inches(0.08), col_w, Inches(0.4),
                      f"{ICON_DIAMOND}  Points cles", font_size=14, bold=True,
                      color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

        y = col_y + Inches(0.7)
        for item in left_items:
            clean = item.lstrip("- ").strip()
            cls._icon_badge(slide, left_x + Inches(0.3), y + Inches(0.02),
                            Inches(0.28), sc, ICON_CHECK, icon_size=10)
            cls._text_box(slide, left_x + Inches(0.7), y,
                          col_w - Inches(1.0), Inches(0.45),
                          clean, font_size=12, color=COLOR_DARK_TEXT)
            y += Inches(0.55)

        # Right column
        right_sc, right_sc_light = SECTION_COLORS[(section_idx + 1) % len(SECTION_COLORS)]
        if right_items:
            cls._rounded_rect(slide, right_x, col_y, col_w, col_h,
                              right_sc_light, border_color=right_sc)
            cls._rect(slide, right_x, col_y, col_w, Inches(0.5), right_sc)
            cls._text_box(slide, right_x, col_y + Inches(0.08), col_w, Inches(0.4),
                          f"{ICON_STAR}  Avantages", font_size=14, bold=True,
                          color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

            y = col_y + Inches(0.7)
            for item in right_items:
                clean = item.lstrip("- ").strip()
                cls._icon_badge(slide, right_x + Inches(0.3), y + Inches(0.02),
                                Inches(0.28), right_sc, ICON_STAR, icon_size=10)
                cls._text_box(slide, right_x + Inches(0.7), y,
                              col_w - Inches(1.0), Inches(0.45),
                              clean, font_size=12, color=COLOR_DARK_TEXT)
                y += Inches(0.55)

        # VS / separator between columns
        vs_y = col_y + col_h / 2 - Inches(0.3)
        cls._icon_badge(slide, Inches(6.25), vs_y, Inches(0.5),
                        COLOR_PRIMARY, ICON_ARROW_R, icon_size=16)

        cls._add_notes(slide, speaker_notes)

    @classmethod
    def _create_quote_slide(cls, prs, title, items, subtitle="",
                             speaker_notes="", section_label="",
                             section_idx=0):
        """Quote/highlight slide for key messages or value propositions."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sc, sc_light = SECTION_COLORS[section_idx % len(SECTION_COLORS)]

        cls._add_background(slide, COLOR_PRIMARY)

        # Decorative circles
        cls._oval(slide, cls.SLIDE_W - Inches(5), Inches(-2),
                  Inches(7), Inches(7), COLOR_SECONDARY)
        cls._oval(slide, Inches(-3), cls.SLIDE_H - Inches(3),
                  Inches(5), Inches(5), COLOR_SECONDARY)
        cls._oval(slide, cls.SLIDE_W - Inches(2), cls.SLIDE_H - Inches(2),
                  Inches(3), Inches(3), sc)

        # Decorative dots
        cls._add_decorative_dots(slide, Inches(1), Inches(0.5), cols=4, rows=3,
                                 color=COLOR_ACCENT_LIGHT, spacing=Inches(0.25))
        cls._add_decorative_dots(slide, cls.SLIDE_W - Inches(2.5),
                                 cls.SLIDE_H - Inches(2), cols=4, rows=3,
                                 color=COLOR_ACCENT_LIGHT, spacing=Inches(0.25))

        # Section label at top
        if section_label:
            cls._text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.4),
                          section_label.upper(), font_size=11, bold=True,
                          color=sc_light, alignment=PP_ALIGN.CENTER)

        # Large opening quote mark
        cls._text_box(slide, Inches(1), Inches(1.0), Inches(2), Inches(1.5),
                      ICON_QUOTE_L, font_size=72, bold=True,
                      color=sc, alignment=PP_ALIGN.LEFT)

        # Main quote text (first item or concatenated)
        if items:
            main_text = items[0].lstrip("- ").strip()
            cls._text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(2.5),
                          main_text, font_size=26, bold=True,
                          color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

        # Closing quote
        cls._text_box(slide, cls.SLIDE_W - Inches(3), Inches(4.5),
                      Inches(2), Inches(1.5),
                      ICON_QUOTE_R, font_size=72, bold=True,
                      color=sc, alignment=PP_ALIGN.RIGHT)

        # Additional items as supporting points below
        if len(items) > 1:
            y = Inches(5.2)
            for item in items[1:4]:
                clean = item.lstrip("- ").strip()
                cls._text_box(slide, Inches(2), y, Inches(9), Inches(0.4),
                              f"{ICON_ARROW_R}  {clean}", font_size=13,
                              color=COLOR_ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
                y += Inches(0.45)

        # Accent bar
        cls._rect(slide, Inches(4), Inches(4.9), Inches(5), Inches(0.04), sc)

        # Footer
        cls._rect(slide, Inches(0), cls.SLIDE_H - Inches(0.3),
                  cls.SLIDE_W, Inches(0.3), COLOR_SECONDARY)

        cls._add_notes(slide, speaker_notes)

    @classmethod
    def _create_pyramid_slide(cls, prs, title, items, subtitle="",
                               speaker_notes="", section_label="",
                               section_idx=0):
        """Pyramid/funnel visualization - items ordered from top (narrow) to bottom (wide)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_WHITE)

        sc, sc_light = SECTION_COLORS[section_idx % len(SECTION_COLORS)]
        icon = _get_section_icon(section_label or title)

        cls._add_left_sidebar(slide, accent_color=sc, icon=icon)
        cls._add_slide_header(slide, title, subtitle, section_label, accent_color=sc)
        cls._add_footer_bar(slide, section_color=sc)

        n_items = min(len(items), 6)
        if n_items == 0:
            cls._add_notes(slide, speaker_notes)
            return

        # Build pyramid layers (top = narrowest, bottom = widest)
        center_x = Inches(6.5)
        start_y = Inches(1.5)
        min_w = Inches(3.5)
        max_w = Inches(11.0)
        layer_h = Inches(0.85)
        gap = Inches(0.1)

        for i, item_text in enumerate(items[:n_items]):
            clean = item_text.lstrip("- ").strip()
            item_sc, item_sc_light = SECTION_COLORS[(section_idx + i) % len(SECTION_COLORS)]

            # Width increases from top to bottom
            progress = i / max(n_items - 1, 1)
            w = min_w + (max_w - min_w) * progress
            x = center_x - w / 2
            y = start_y + i * (layer_h + gap)

            # Layer rectangle with rounded corners
            cls._rounded_rect(slide, x, y, w, layer_h, item_sc)

            # Icon on the left side of the layer
            layer_icon = [ICON_TROPHY, ICON_ROCKET, ICON_LIGHT, ICON_GEAR,
                          ICON_SHIELD, ICON_CHART][i % 6]
            cls._text_box(slide, x + Inches(0.3), y + Inches(0.15),
                          Inches(0.5), Inches(0.5),
                          layer_icon, font_size=18, color=COLOR_WHITE,
                          alignment=PP_ALIGN.CENTER)

            # Text
            cls._text_box(slide, x + Inches(0.9), y + Inches(0.15),
                          w - Inches(1.2), layer_h - Inches(0.3),
                          clean, font_size=13, bold=True,
                          color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

        cls._add_notes(slide, speaker_notes)

    @classmethod
    def _create_dashboard_slide(cls, prs, title, items, subtitle="",
                                 speaker_notes="", section_label="",
                                 section_idx=0):
        """Dashboard-style slide with mixed metrics and text cards."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_LIGHT_BG)

        sc, sc_light = SECTION_COLORS[section_idx % len(SECTION_COLORS)]
        icon = _get_section_icon(section_label or title)

        cls._add_left_sidebar(slide, accent_color=sc, icon=icon)
        cls._add_slide_header(slide, title, subtitle, section_label, accent_color=sc)
        cls._add_footer_bar(slide, section_color=sc)

        n_items = min(len(items), 6)
        if n_items == 0:
            cls._add_notes(slide, speaker_notes)
            return

        # Layout: top row of small metric cards + bottom larger info cards
        top_items = items[:min(3, n_items)]
        bottom_items = items[min(3, n_items):n_items]

        # Top row - small metric-style cards
        top_card_w = Inches(3.6)
        top_card_h = Inches(2.2)
        top_start_x = Inches(0.8)
        top_y = Inches(1.5)
        top_gap = Inches(0.3)

        dash_icons = [ICON_CHART, ICON_ROCKET, ICON_TROPHY,
                      ICON_GEAR, ICON_SHIELD, ICON_LIGHT]

        for i, item_text in enumerate(top_items):
            clean = item_text.lstrip("- ").strip()
            item_sc, item_sc_light = SECTION_COLORS[(section_idx + i) % len(SECTION_COLORS)]
            x = top_start_x + i * (top_card_w + top_gap)

            # Card
            cls._rounded_rect(slide, x, top_y, top_card_w, top_card_h, COLOR_WHITE)
            cls._rect(slide, x, top_y, top_card_w, Inches(0.05), item_sc)

            # Large icon
            icon_text = dash_icons[i % len(dash_icons)]
            cls._icon_badge(slide, x + Inches(0.3), top_y + Inches(0.25),
                            Inches(0.55), item_sc, icon_text, icon_size=18)

            # Text
            cls._text_box(slide, x + Inches(1.0), top_y + Inches(0.2),
                          top_card_w - Inches(1.3), top_card_h - Inches(0.4),
                          clean, font_size=12, color=COLOR_DARK_TEXT)

        # Bottom row - wider cards
        if bottom_items:
            bot_n = len(bottom_items)
            bot_total_w = Inches(11.3)
            bot_card_w = bot_total_w / bot_n - Inches(0.2)
            bot_card_h = Inches(2.5)
            bot_start_x = Inches(0.8)
            bot_y = top_y + top_card_h + Inches(0.3)

            for i, item_text in enumerate(bottom_items):
                clean = item_text.lstrip("- ").strip()
                item_sc, item_sc_light = SECTION_COLORS[(section_idx + 3 + i) % len(SECTION_COLORS)]
                x = bot_start_x + i * (bot_card_w + Inches(0.2))

                cls._rounded_rect(slide, x, bot_y, bot_card_w, bot_card_h,
                                  COLOR_WHITE, border_color=item_sc_light)
                # Left colored accent
                cls._rect(slide, x, bot_y + Inches(0.15), Inches(0.05),
                          bot_card_h - Inches(0.3), item_sc)

                icon_text = dash_icons[(3 + i) % len(dash_icons)]
                cls._icon_badge(slide, x + Inches(0.25), bot_y + Inches(0.25),
                                Inches(0.45), item_sc, icon_text, icon_size=14)

                cls._text_box(slide, x + Inches(0.85), bot_y + Inches(0.15),
                              bot_card_w - Inches(1.1), bot_card_h - Inches(0.3),
                              clean, font_size=12, color=COLOR_DARK_TEXT)

        cls._add_notes(slide, speaker_notes)

    @classmethod
    def _create_visual_content_slide(cls, prs, title, bullet_points, subtitle="",
                                      speaker_notes="", section_label="",
                                      section_idx=0, layout=LAYOUT_BULLETS):
        """Route to the appropriate visual layout based on the layout field."""
        # Normalize all bullet items to plain strings (AI may return dicts)
        bullet_points = [_normalize_bullet(b) for b in bullet_points]
        if layout == LAYOUT_PROCESS:
            cls._create_process_flow_slide(prs, title, bullet_points, subtitle,
                                           speaker_notes, section_label, section_idx)
        elif layout == LAYOUT_TIMELINE:
            cls._create_timeline_slide(prs, title, bullet_points, subtitle,
                                       speaker_notes, section_label, section_idx)
        elif layout == LAYOUT_ICON_GRID:
            cls._create_icon_grid_slide(prs, title, bullet_points, subtitle,
                                        speaker_notes, section_label, section_idx)
        elif layout == LAYOUT_COMPARISON:
            cls._create_comparison_slide(prs, title, bullet_points, subtitle,
                                         speaker_notes, section_label, section_idx)
        elif layout == LAYOUT_QUOTE:
            cls._create_quote_slide(prs, title, bullet_points, subtitle,
                                    speaker_notes, section_label, section_idx)
        elif layout == LAYOUT_PYRAMID:
            cls._create_pyramid_slide(prs, title, bullet_points, subtitle,
                                      speaker_notes, section_label, section_idx)
        elif layout == LAYOUT_DASHBOARD:
            cls._create_dashboard_slide(prs, title, bullet_points, subtitle,
                                        speaker_notes, section_label, section_idx)
        else:
            cls._create_content_slide(prs, title, bullet_points, subtitle,
                                      speaker_notes, section_label, section_idx)

    @classmethod
    def _create_key_figures_slide(cls, prs, title, figures):
        """Key figures with large numbers, color-coded cards, and icon badges."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_WHITE)
        cls._add_slide_header(slide, title)
        cls._add_footer_bar(slide)

        # Decorative background shapes
        cls._oval(slide, cls.SLIDE_W - Inches(3), Inches(1.5),
                  Inches(4), Inches(4), COLOR_LIGHT_BG)

        cols = min(len(figures), 4)
        box_width = Inches(11) / cols
        start_x = Inches(1)

        fig_icons = [ICON_TARGET, ICON_STAR, ICON_DIAMOND, ICON_CHECK,
                     ICON_CIRCLE, ICON_SQUARE, ICON_RING, ICON_TRIANGLE]

        for i, fig in enumerate(figures):
            col = i % cols
            row = i // cols
            x = start_x + col * box_width
            y = Inches(1.8) + row * Inches(2.5)
            sc, sc_light = SECTION_COLORS[i % len(SECTION_COLORS)]

            card_w = box_width - Inches(0.4)
            # Card with light background
            cls._rounded_rect(slide, x, y, card_w, Inches(2.2), sc_light)
            # Top color bar
            cls._rect(slide, x, y, card_w, Inches(0.06), sc)

            # Icon badge top-right
            cls._icon_badge(slide, x + card_w - Inches(0.5), y + Inches(0.15),
                            Inches(0.35), sc, fig_icons[i % len(fig_icons)], icon_size=12)

            # Value
            cls._text_box(slide, x, y + Inches(0.3), card_w, Inches(0.9),
                          fig.get("value", ""), font_size=40, bold=True,
                          color=sc, alignment=PP_ALIGN.CENTER)

            # Separator
            cls._rect(slide, x + Inches(0.4), y + Inches(1.3),
                      card_w - Inches(0.8), Inches(0.02), sc)

            # Label
            cls._text_box(slide, x + Inches(0.15), y + Inches(1.4),
                          card_w - Inches(0.3), Inches(0.7),
                          fig.get("label", ""), font_size=12,
                          color=COLOR_MUTED, alignment=PP_ALIGN.CENTER)

    @classmethod
    def _create_strengths_slide(cls, prs, strengths):
        """Strengths slides with numbered badges and star icons."""
        strengths = [_normalize_bullet(s) for s in strengths]
        per_slide = 5
        for page in range(0, len(strengths), per_slide):
            chunk = strengths[page:page + per_slide]
            slide_num = page // per_slide + 1
            total_pages = math.ceil(len(strengths) / per_slide)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            cls._add_background(slide, COLOR_LIGHT_BG)

            title_sfx = f" ({slide_num}/{total_pages})" if total_pages > 1 else ""
            cls._add_slide_header(slide, f"NOS FORCES{title_sfx}", accent_color=COLOR_WARN)
            cls._add_footer_bar(slide, section_color=COLOR_WARN)

            # Decorative star pattern
            cls._add_decorative_dots(slide, cls.SLIDE_W - Inches(2.5), Inches(1.5),
                                     cols=4, rows=3, color=COLOR_WARN_LIGHT, spacing=Inches(0.3))

            y = Inches(1.4)
            for i, strength in enumerate(chunk):
                global_idx = page + i
                sc, sc_light = SECTION_COLORS[global_idx % len(SECTION_COLORS)]

                # Card with shadow effect (offset darker rect)
                cls._rounded_rect(slide, Inches(0.55), y + Inches(0.03),
                                  Inches(12), Inches(0.85), COLOR_ACCENT_PALE)
                cls._rounded_rect(slide, Inches(0.5), y, Inches(12), Inches(0.85), COLOR_WHITE)

                # Number badge
                cls._icon_badge(slide, Inches(0.7), y + Inches(0.12),
                                Inches(0.6), sc, str(global_idx + 1), icon_size=18)

                # Star icon
                cls._icon_badge(slide, Inches(11.3), y + Inches(0.17),
                                Inches(0.5), COLOR_WARN, ICON_STAR, icon_size=14)

                # Text
                cls._text_box(slide, Inches(1.5), y + Inches(0.15),
                              Inches(9.5), Inches(0.55), strength,
                              font_size=14, color=COLOR_DARK_TEXT)

                y += Inches(1.0)

    @classmethod
    def _create_qa_slide(cls, prs, questions):
        """Q&A preparation slides."""
        per_slide = 3
        for page in range(0, len(questions), per_slide):
            chunk = questions[page:page + per_slide]

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            cls._add_background(slide, COLOR_WHITE)
            cls._add_slide_header(slide, "QUESTIONS ANTICIPEES", accent_color=COLOR_PURPLE)
            cls._add_footer_bar(slide, section_color=COLOR_PURPLE)

            y = Inches(1.4)
            for i, qa in enumerate(chunk):
                q = qa.get("question", "")
                a = qa.get("answer", "")

                # Question card
                cls._rounded_rect(slide, Inches(0.5), y, Inches(12), Inches(1.6),
                                  COLOR_PURPLE_LIGHT, border_color=COLOR_PURPLE)

                # Q badge
                cls._icon_badge(slide, Inches(0.7), y + Inches(0.15),
                                Inches(0.45), COLOR_PURPLE, "Q", icon_size=16)

                # Question text
                cls._text_box(slide, Inches(1.3), y + Inches(0.1),
                              Inches(10.5), Inches(0.5), q,
                              font_size=13, bold=True, color=COLOR_PRIMARY)

                # Answer with arrow
                if a:
                    display = a[:280] + "..." if len(a) > 280 else a
                    cls._icon_badge(slide, Inches(1.3), y + Inches(0.7),
                                    Inches(0.25), COLOR_SUCCESS, ICON_CHECK, icon_size=9)
                    cls._text_box(slide, Inches(1.7), y + Inches(0.65),
                                  Inches(10.2), Inches(0.85), display,
                                  font_size=11, color=COLOR_MUTED)

                y += Inches(1.8)

    @classmethod
    def _create_summary_slide(cls, prs, project_name, strengths):
        """Visual recap slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_PRIMARY)

        # Decorative
        cls._oval(slide, cls.SLIDE_W - Inches(4), Inches(-2),
                  Inches(6), Inches(6), COLOR_SECONDARY)
        cls._oval(slide, Inches(-2), cls.SLIDE_H - Inches(3),
                  Inches(4), Inches(4), COLOR_SECONDARY)

        # Dot pattern
        cls._add_decorative_dots(slide, Inches(9), Inches(0.5), cols=5, rows=3,
                                 color=COLOR_ACCENT_LIGHT, spacing=Inches(0.3))

        cls._text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                      "EN RESUME", font_size=34, bold=True, color=COLOR_WHITE)

        cls._rect(slide, Inches(0.8), Inches(1.3), Inches(3.5), Inches(0.04), COLOR_ACCENT)

        y = Inches(1.8)
        items = strengths[:6] if strengths else []
        for i, item in enumerate(items):
            sc, _ = SECTION_COLORS[i % len(SECTION_COLORS)]
            cls._icon_badge(slide, Inches(0.8), y + Inches(0.02),
                            Inches(0.3), sc, ICON_CHECK, icon_size=11)
            cls._text_box(slide, Inches(1.3), y, Inches(10), Inches(0.5),
                          item, font_size=15, color=COLOR_WHITE)
            y += Inches(0.6)

        cls._add_footer_bar(slide)

    @classmethod
    def _create_closing_slide(cls, prs, company_name, client_name):
        """Professional closing slide with rich visual design."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cls._add_background(slide, COLOR_PRIMARY)

        # Rich decorative circles (multiple layers)
        cls._oval(slide, Inches(-3), Inches(-3),
                  Inches(8), Inches(8), COLOR_SECONDARY)
        cls._oval(slide, cls.SLIDE_W - Inches(5), cls.SLIDE_H - Inches(5),
                  Inches(8), Inches(8), COLOR_SECONDARY)
        cls._oval(slide, Inches(4), Inches(-1), Inches(3), Inches(3), COLOR_ACCENT)
        cls._oval(slide, cls.SLIDE_W - Inches(2), Inches(1),
                  Inches(2), Inches(2), COLOR_ACCENT)
        cls._oval(slide, Inches(1), cls.SLIDE_H - Inches(3),
                  Inches(2.5), Inches(2.5), COLOR_ACCENT)

        # Dot patterns (multiple clusters)
        cls._add_decorative_dots(slide, Inches(0.5), Inches(0.5), cols=4, rows=3,
                                 color=COLOR_ACCENT_LIGHT, spacing=Inches(0.25))
        cls._add_decorative_dots(slide, cls.SLIDE_W - Inches(2.5),
                                 cls.SLIDE_H - Inches(2), cols=5, rows=3,
                                 color=COLOR_ACCENT_LIGHT, spacing=Inches(0.25))

        # Top accent bars
        cls._rect(slide, Inches(0), Inches(0), cls.SLIDE_W, Inches(0.06), COLOR_ACCENT)
        cls._rect(slide, Inches(0), Inches(0.06), cls.SLIDE_W, Inches(0.03), COLOR_TEAL)

        # Handshake icon
        cls._text_box(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.8),
                      ICON_HANDSHAKE, font_size=40, color=COLOR_ACCENT,
                      alignment=PP_ALIGN.CENTER)

        cls._text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
                      "MERCI", font_size=64, bold=True,
                      color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

        # Double accent bars
        cls._rect(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(0.05), COLOR_ACCENT)
        cls._rect(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.03), COLOR_TEAL)

        cls._text_box(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.6),
                      "Questions & Echanges", font_size=30,
                      color=COLOR_ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

        # Contact info with icon badges
        info_y = Inches(5.0)
        if company_name:
            cls._icon_badge(slide, Inches(3.5), info_y, Inches(0.35),
                            COLOR_TEAL, ICON_DIAMOND, icon_size=12)
            cls._text_box(slide, Inches(4.0), info_y, Inches(5), Inches(0.4),
                          company_name, font_size=16, bold=True,
                          color=COLOR_WHITE, alignment=PP_ALIGN.LEFT)
            info_y += Inches(0.5)
        if client_name:
            cls._icon_badge(slide, Inches(3.5), info_y, Inches(0.35),
                            COLOR_ACCENT, ICON_TARGET, icon_size=12)
            cls._text_box(slide, Inches(4.0), info_y, Inches(5), Inches(0.4),
                          f"Client : {client_name}", font_size=14,
                          color=COLOR_ACCENT_LIGHT, alignment=PP_ALIGN.LEFT)

        # Bottom bar (dual-tone)
        cls._rect(slide, Inches(0), cls.SLIDE_H - Inches(0.55),
                  cls.SLIDE_W, Inches(0.55), COLOR_SECONDARY)
        cls._rect(slide, Inches(0), cls.SLIDE_H - Inches(0.55),
                  cls.SLIDE_W, Inches(0.04), COLOR_ACCENT)

    # ── Main generation ──

    @classmethod
    def generate_presentation(cls, project_name, client_name, company_name,
                              rfp_reference, soutenance_data):
        """Generate a complete soutenance PowerPoint from structured AI data."""
        prs = Presentation()
        prs.slide_width = cls.SLIDE_W
        prs.slide_height = cls.SLIDE_H

        sections = soutenance_data.get("sections", [])
        key_figures = soutenance_data.get("key_figures", [])
        strengths = soutenance_data.get("strengths", [])
        script = soutenance_data.get("script", {})
        qa_prep = script.get("qa_preparation", {})
        total_sections = len(sections)

        # 1. Cover slide
        cls._create_title_slide(prs, project_name, client_name, company_name, rfp_reference)

        # 2. Agenda slide
        cls._create_agenda_slide(prs, sections)

        # 3. Content slides per section
        for sec_idx, section in enumerate(sections):
            section_title = section.get("title", "")
            duration = section.get("duration", "")

            # Section divider
            cls._create_section_divider(prs, sec_idx + 1, section_title,
                                        duration, total_sections)

            # Section content slides
            for slide_data in section.get("slides", []):
                layout = slide_data.get("layout", LAYOUT_BULLETS)
                cls._create_visual_content_slide(
                    prs,
                    title=slide_data.get("title", ""),
                    bullet_points=slide_data.get("bullets", []),
                    subtitle=slide_data.get("subtitle", ""),
                    speaker_notes=slide_data.get("speaker_notes", ""),
                    section_label=section_title,
                    section_idx=sec_idx,
                    layout=layout,
                )

        # 4. Key figures slide
        if key_figures:
            cls._create_key_figures_slide(prs, "CHIFFRES CLES", key_figures)

        # 5. Strengths slides
        if strengths:
            cls._create_strengths_slide(prs, strengths)

        # 6. Q&A preparation slides
        expected_qs = qa_prep.get("expected_questions", [])
        if expected_qs:
            cls._create_qa_slide(prs, expected_qs)

        # 7. Summary slide
        cls._create_summary_slide(prs, project_name, strengths)

        # 8. Closing slide
        cls._create_closing_slide(prs, company_name, client_name)

        # Save
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream
